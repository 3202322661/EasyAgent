import os
import tempfile
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ────────────────────────── 常量映射表 ──────────────────────────

COLOR_MAP = {
    "red": RGBColor(0xFF, 0x00, 0x00),
    "green": RGBColor(0x00, 0x80, 0x00),
    "blue": RGBColor(0x00, 0x00, 0xFF),
    "black": RGBColor(0x00, 0x00, 0x00),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "yellow": RGBColor(0xFF, 0xFF, 0x00),
    "orange": RGBColor(0xFF, 0xA5, 0x00),
    "purple": RGBColor(0x80, 0x00, 0x80),
    "gray": RGBColor(0x80, 0x80, 0x80),
    "dark_red": RGBColor(0x8B, 0x00, 0x00),
    "dark_green": RGBColor(0x00, 0x64, 0x00),
    "dark_blue": RGBColor(0x00, 0x00, 0x8B),
    "light_gray": RGBColor(0xD3, 0xD3, 0xD3),
    "cyan": RGBColor(0x00, 0xFF, 0xFF),
    "magenta": RGBColor(0xFF, 0x00, 0xFF),
    "brown": RGBColor(0xA5, 0x2A, 0x2A),
    "navy": RGBColor(0x00, 0x00, 0x80),
    "teal": RGBColor(0x00, 0x80, 0x80),
}

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}

LAYOUT_NAMES = {
    "title": 0,           # 标题幻灯片
    "title_content": 1,   # 标题和内容
    "section": 2,         # 节标题
    "two_content": 3,     # 两栏内容
    "blank": 6,           # 空白
    "content_caption": 8, # 内容和标题
    "picture_caption": 10,# 图片和标题
}

SHAPE_TYPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "right_triangle": MSO_SHAPE.RIGHT_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "pentagon": MSO_SHAPE.PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "arrow_right": MSO_SHAPE.RIGHT_ARROW,
    "arrow_left": MSO_SHAPE.LEFT_ARROW,
    "arrow_up": MSO_SHAPE.UP_ARROW,
    "arrow_down": MSO_SHAPE.DOWN_ARROW,
    "star": MSO_SHAPE.STAR_5_POINT,
    "star_5": MSO_SHAPE.STAR_5_POINT,
    "heart": MSO_SHAPE.HEART,
    "cloud": MSO_SHAPE.CLOUD,
    "sun": MSO_SHAPE.SUN,
    "moon": MSO_SHAPE.MOON,
    "lightning": MSO_SHAPE.LIGHTNING_BOLT,
    "cross": MSO_SHAPE.CROSS,
    
    
    
    
    
    "cube": MSO_SHAPE.CUBE,
    
    
    "no_symbol": MSO_SHAPE.NO_SYMBOL,
    "chevron": MSO_SHAPE.CHEVRON,
    "flowchart_process": MSO_SHAPE.FLOWCHART_PROCESS,
    "flowchart_decision": MSO_SHAPE.FLOWCHART_DECISION,
    "flowchart_data": MSO_SHAPE.FLOWCHART_DATA,
    "flowchart_start": MSO_SHAPE.FLOWCHART_TERMINATOR,
}

# ────────────────────────── 内部辅助函数 ──────────────────────────

def _parse_color(color_val):
    """解析颜色值，支持英文名称和十六进制格式"""
    if color_val is None:
        return None
    color_lower = color_val.lower().replace(" ", "_")
    if color_lower in COLOR_MAP:
        return COLOR_MAP[color_lower]
    if color_val.startswith("#"):
        hex_color = color_val[1:]
        if len(hex_color) == 6:
            try:
                return RGBColor(
                    int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16),
                )
            except ValueError:
                pass
    if len(color_val) == 6 and all(c in "0123456789abcdefABCDEF" for c in color_val):
        try:
            return RGBColor(
                int(color_val[0:2], 16),
                int(color_val[2:4], 16),
                int(color_val[4:6], 16),
            )
        except ValueError:
            pass
    return None


def _parse_align(align_str):
    """解析对齐方式"""
    if align_str is None:
        return None
    return ALIGN_MAP.get(align_str.lower())


def _safe_open_ppt(file_path):
    """安全地打开 PPT 文件，如果不存在则创建新文档"""
    if not os.path.exists(file_path):
        return Presentation()
    try:
        return Presentation(file_path)
    except Exception:
        return Presentation()


def _safe_save_ppt(prs, file_path):
    """安全地保存 PPT 文件，先保存到临时文件再替换"""
    try:
        fd, tmp_path = tempfile.mkstemp(suffix='.pptx')
        os.close(fd)
        prs.save(tmp_path)
        shutil.move(tmp_path, file_path)
        return True
    except Exception:
        return False


def _get_layout(prs, layout_name_or_index):
    """根据布局名称或索引获取幻灯片布局"""
    if layout_name_or_index is None:
        return prs.slide_layouts[6]  # 默认使用空白布局

    if isinstance(layout_name_or_index, str):
        layout_key = layout_name_or_index.lower().replace(" ", "_")
        idx = LAYOUT_NAMES.get(layout_key)
        if idx is not None:
            try:
                return prs.slide_layouts[idx]
            except IndexError:
                pass
        return prs.slide_layouts[6]

    if isinstance(layout_name_or_index, int):
        try:
            return prs.slide_layouts[layout_name_or_index]
        except IndexError:
            return prs.slide_layouts[6]

    return prs.slide_layouts[6]


def _set_shape_fill(shape, color=None, transparency=None):
    """设置形状填充颜色"""
    if color is None:
        return
    parsed = _parse_color(color)
    if parsed is None:
        return
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = parsed
        if transparency is not None:
            from pptx.oxml.ns import qn
            solidFill = shape.fill._fill
            srgb = solidFill.find(qn('a:solidFill'))
            if srgb is not None:
                clr = srgb[0] if len(srgb) else None
                if clr is not None:
                    clr.set('alpha', str(int(transparency * 1000)))
    except Exception:
        pass


def _set_shape_line(shape, color=None, width=None):
    """设置形状边框"""
    if color:
        parsed = _parse_color(color)
        if parsed:
            try:
                shape.line.color.rgb = parsed
            except Exception:
                pass
    if width is not None:
        try:
            shape.line.width = Pt(width)
        except Exception:
            pass


def _set_text_format(text_frame, text, font_name=None, font_size=None,
                     bold=None, italic=None, color=None, alignment=None):
    """设置文本框内容的格式"""
    text_frame.text = str(text)
    para = text_frame.paragraphs[0]

    if alignment:
        parsed_align = _parse_align(alignment)
        if parsed_align:
            para.alignment = parsed_align

    for run in para.runs:
        if font_name:
            try:
                run.font.name = font_name
            except Exception:
                pass
        if font_size is not None:
            try:
                run.font.size = Pt(font_size)
            except Exception:
                pass
        if bold is not None:
            run.font.bold = bold
        if italic is not None:
            run.font.italic = italic
        if color:
            parsed = _parse_color(color)
            if parsed:
                try:
                    run.font.color.rgb = parsed
                except Exception:
                    pass


def _add_slide_master_decorator(prs):
    """确保幻灯片母版存在（内部辅助）"""
    pass


def _set_slide_bg_color(slide, color):
    """设置幻灯片背景颜色"""
    parsed = _parse_color(color)
    if parsed is None:
        return
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = parsed
    except Exception:
        pass

# ────────────────────────── 公开接口函数 ──────────────────────────

def create_ppt(file_path: str, title: str = None) -> str:
    """
    创建一个新的空白 PPT 文件。

    Args:
        file_path: 保存路径（.pptx 格式）
        title: 可选的标题幻灯片标题文本

    Returns:
        操作结果字符串
    """
    try:
        prs = Presentation()
        if title:
            slide_layout = prs.slide_layouts[0]  # 标题幻灯片
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            slide_count = len(prs.slides)
            return f"✓ PPT 文件已创建: '{file_path}'（含 {slide_count} 张幻灯片）"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 创建 PPT 失败: {str(e)}"


def add_slide(file_path: str, layout: str = "blank",
              title: str = None, content: str = None) -> str:
    """
    向 PPT 文件添加一张新幻灯片。

    Args:
        file_path: PPT 文件路径
        layout: 布局名称（title/title_content/section/two_content/blank/
                content_caption/picture_caption）或索引数字
        title: 幻灯片标题文本（仅对含标题占位符的布局有效）
        content: 幻灯片内容文本（仅对含内容占位符的布局有效）

    Returns:
        操作结果字符串
    """
    try:
        prs = _safe_open_ppt(file_path)
        slide_layout = _get_layout(prs, layout)
        slide = prs.slides.add_slide(slide_layout)

        # 设置标题
        if title and slide.shapes.title:
            slide.shapes.title.text = title

        # 设置内容（如果有内容占位符）
        if content:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    shape.text = content
                    break

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            total = len(prs.slides)
            return f"✓ 已添加幻灯片（使用 '{layout}' 布局），共 {total} 张幻灯片"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 添加幻灯片失败: {str(e)}"


def add_textbox(file_path: str, slide_index: int = -1,
                left: float = 1.0, top: float = 1.0,
                width: float = 8.0, height: float = 1.5,
                text: str = "", font_name: str = None,
                font_size: int = 18, bold: bool = None,
                italic: bool = None, color: str = None,
                alignment: str = None) -> str:
    """
    在指定幻灯片中添加文本框。

    Args:
        file_path: PPT 文件路径
        slide_index: 幻灯片索引（从 1 开始，-1 表示最后一张）
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 宽度（英寸）
        height: 高度（英寸）
        text: 文本框内容
        font_name: 字体名称
        font_size: 字号（磅）
        bold: 是否加粗
        italic: 是否斜体
        color: 文字颜色
        alignment: 对齐方式（left/center/right/justify）

    Returns:
        操作结果字符串
    """
    if not os.path.exists(file_path):
        return f"[ERROR] 文件 '{file_path}' 不存在。"

    try:
        prs = Presentation(file_path)
        slides = prs.slides
        if not slides:
            return "[ERROR] PPT 中没有幻灯片，请先添加幻灯片。"

        if slide_index == -1:
            slide = slides[-1]
        else:
            if slide_index < 1 or slide_index > len(slides):
                return f"[ERROR] 幻灯片索引 {slide_index} 超出范围 (1-{len(slides)})。"
            slide = slides[slide_index - 1]

        left_inches = Inches(left)
        top_inches = Inches(top)
        width_inches = Inches(width)
        height_inches = Inches(height)

        txBox = slide.shapes.add_textbox(left_inches, top_inches,
                                          width_inches, height_inches)
        tf = txBox.text_frame
        tf.word_wrap = True

        _set_text_format(tf, text, font_name=font_name, font_size=font_size,
                         bold=bold, italic=italic, color=color,
                         alignment=alignment)

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            return f"✓ 文本框已添加到第 {slide_index if slide_index != -1 else '最后'} 张幻灯片"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 添加文本框失败: {str(e)}"


def add_shape(file_path: str, slide_index: int = -1,
              shape_type: str = "rectangle",
              left: float = 1.0, top: float = 1.0,
              width: float = 3.0, height: float = 2.0,
              fill_color: str = None, line_color: str = None,
              line_width: float = None, text: str = None,
              font_name: str = None, font_size: int = None,
              font_color: str = None, bold: bool = None,
              alignment: str = None) -> str:
    """
    在指定幻灯片中添加形状（矩形、圆形、箭头等）。

    Args:
        file_path: PPT 文件路径
        slide_index: 幻灯片索引（从 1 开始，-1 表示最后一张）
        shape_type: 形状类型（rectangle/oval/triangle/arrow_right/star/heart等）
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 宽度（英寸）
        height: 高度（英寸）
        fill_color: 填充颜色
        line_color: 边框颜色
        line_width: 边框宽度（磅）
        text: 形状内文字
        font_name: 字体名称
        font_size: 字号（磅）
        font_color: 文字颜色
        bold: 是否加粗
        alignment: 对齐方式

    Returns:
        操作结果字符串
    """
    if not os.path.exists(file_path):
        return f"[ERROR] 文件 '{file_path}' 不存在。"

    try:
        prs = Presentation(file_path)
        slides = prs.slides
        if not slides:
            return "[ERROR] PPT 中没有幻灯片，请先添加幻灯片。"

        if slide_index == -1:
            slide = slides[-1]
        else:
            if slide_index < 1 or slide_index > len(slides):
                return f"[ERROR] 幻灯片索引 {slide_index} 超出范围 (1-{len(slides)})。"
            slide = slides[slide_index - 1]

        # 解析形状类型
        shape_enum = SHAPE_TYPE_MAP.get(shape_type.lower().replace(" ", "_"))
        if shape_enum is None:
            shape_enum = MSO_SHAPE.RECTANGLE

        shape = slide.shapes.add_shape(
            shape_enum,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )

        # 设置填充
        _set_shape_fill(shape, color=fill_color)

        # 设置边框
        _set_shape_line(shape, color=line_color, width=line_width)

        # 设置文字
        if text is not None and hasattr(shape, 'text_frame'):
            _set_text_format(shape.text_frame, text, font_name=font_name,
                             font_size=font_size, color=font_color,
                             bold=bold, alignment=alignment)

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            return f"✓ 形状（{shape_type}）已添加到第 {slide_index if slide_index != -1 else '最后'} 张幻灯片"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 添加形状失败: {str(e)}"


def add_table_slide(file_path: str, slide_index: int = -1,
                    data: list = None, headers: list = None,
                    left: float = 1.0, top: float = 2.0,
                    width: float = 8.0, height: float = None,
                    font_name: str = None, font_size: int = 12,
                    header_color: str = None, header_bold: bool = True,
                    alignment: str = None) -> str:
    """
    在指定幻灯片中添加表格。

    Args:
        file_path: PPT 文件路径
        slide_index: 幻灯片索引（从 1 开始，-1 表示最后一张）
        data: 表格数据，二维列表，如 [['A1','B1'],['A2','B2']]
        headers: 表头列表，如 ['列1', '列2']
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 表格总宽度（英寸）
        height: 表格总高度（英寸，不传则自适应）
        font_name: 字体名称
        font_size: 字号（磅）
        header_color: 表头文字颜色
        header_bold: 表头是否加粗
        alignment: 对齐方式

    Returns:
        操作结果字符串
    """
    if not os.path.exists(file_path):
        return f"[ERROR] 文件 '{file_path}' 不存在。"
    if not data:
        return "[ERROR] 表格数据为空。"

    try:
        prs = Presentation(file_path)
        slides = prs.slides
        if not slides:
            return "[ERROR] PPT 中没有幻灯片，请先添加幻灯片。"

        if slide_index == -1:
            slide = slides[-1]
        else:
            if slide_index < 1 or slide_index > len(slides):
                return f"[ERROR] 幻灯片索引 {slide_index} 超出范围 (1-{len(slides)})。"
            slide = slides[slide_index - 1]

        # 表格行列数
        num_rows = len(data) + (1 if headers else 0)
        num_cols = len(data[0]) if data else 0

        if num_cols == 0:
            return "[ERROR] 表格列数为零。"

        # 计算高度（每行约 0.5 英寸）
        if height is None:
            height = num_rows * 0.5

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        table = table_shape.table

        # 写入表头
        if headers:
            for j, header_text in enumerate(headers):
                if j < num_cols:
                    cell = table.cell(0, j)
                    cell.text = str(header_text)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if font_name:
                                run.font.name = font_name
                            if font_size:
                                run.font.size = Pt(font_size)
                            run.font.bold = header_bold
                            if header_color:
                                parsed = _parse_color(header_color)
                                if parsed:
                                    run.font.color.rgb = parsed
                            if alignment:
                                parsed_align = _parse_align(alignment)
                                if parsed_align:
                                    paragraph.alignment = parsed_align

        # 写入数据
        start_row = 1 if headers else 0
        for i, row_data in enumerate(data):
            for j, cell_text in enumerate(row_data):
                if j < num_cols:
                    cell = table.cell(start_row + i, j)
                    cell.text = str(cell_text)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if font_name:
                                run.font.name = font_name
                            if font_size:
                                run.font.size = Pt(font_size)

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            return f"✓ 表格 ({num_rows} 行 x {num_cols} 列) 已添加到幻灯片"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 添加表格失败: {str(e)}"


def add_image(file_path: str, slide_index: int = -1,
              image_path: str = None,
              left: float = 1.0, top: float = 1.0,
              width: float = None, height: float = None) -> str:
    """
    在指定幻灯片中添加图片。

    Args:
        file_path: PPT 文件路径
        slide_index: 幻灯片索引（从 1 开始，-1 表示最后一张）
        image_path: 图片文件路径
        left: 左边距（英寸）
        top: 上边距（英寸）
        width: 图片宽度（英寸，不传则按比例缩放）
        height: 图片高度（英寸，不传则按比例缩放）

    Returns:
        操作结果字符串
    """
    if not os.path.exists(file_path):
        return f"[ERROR] 文件 '{file_path}' 不存在。"
    if not image_path or not os.path.exists(image_path):
        return f"[ERROR] 图片文件 '{image_path}' 不存在。"

    try:
        prs = Presentation(file_path)
        slides = prs.slides
        if not slides:
            return "[ERROR] PPT 中没有幻灯片，请先添加幻灯片。"

        if slide_index == -1:
            slide = slides[-1]
        else:
            if slide_index < 1 or slide_index > len(slides):
                return f"[ERROR] 幻灯片索引 {slide_index} 超出范围 (1-{len(slides)})。"
            slide = slides[slide_index - 1]

        left_inches = Inches(left)
        top_inches = Inches(top)
        width_inches = Inches(width) if width is not None else None
        height_inches = Inches(height) if height is not None else None

        if width_inches and height_inches:
            slide.shapes.add_picture(image_path, left_inches, top_inches,
                                     width_inches, height_inches)
        elif width_inches:
            slide.shapes.add_picture(image_path, left_inches, top_inches,
                                     width=width_inches)
        elif height_inches:
            slide.shapes.add_picture(image_path, left_inches, top_inches,
                                     height=height_inches)
        else:
            slide.shapes.add_picture(image_path, left_inches, top_inches)

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            return f"✓ 图片已添加到第 {slide_index if slide_index != -1 else '最后'} 张幻灯片"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 添加图片失败: {str(e)}"


def set_slide_background(file_path: str, slide_index: int = -1,
                         color: str = None) -> str:
    """
    设置幻灯片背景颜色。

    Args:
        file_path: PPT 文件路径
        slide_index: 幻灯片索引（从 1 开始，-1 表示所有幻灯片）
        color: 背景颜色

    Returns:
        操作结果字符串
    """
    if not os.path.exists(file_path):
        return f"[ERROR] 文件 '{file_path}' 不存在。"

    try:
        prs = Presentation(file_path)
        slides = prs.slides
        if not slides:
            return "[ERROR] PPT 中没有幻灯片。"

        if slide_index == -1:
            target_slides = list(slides)
            desc = "所有幻灯片"
        else:
            if slide_index < 1 or slide_index > len(slides):
                return f"[ERROR] 幻灯片索引 {slide_index} 超出范围 (1-{len(slides)})。"
            target_slides = [slides[slide_index - 1]]
            desc = f"第 {slide_index} 张幻灯片"

        for slide in target_slides:
            _set_slide_bg_color(slide, color)

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            return f"✓ 已设置 {desc} 的背景颜色"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 设置背景失败: {str(e)}"


def add_title_slide(file_path: str, title: str = "标题",
                    subtitle: str = None,
                    title_color: str = None,
                    subtitle_color: str = None,
                    bg_color: str = None) -> str:
    """
    添加一张标题幻灯片（封面）。

    Args:
        file_path: PPT 文件路径
        title: 主标题文本
        subtitle: 副标题文本（可选）
        title_color: 标题颜色
        subtitle_color: 副标题颜色
        bg_color: 背景颜色

    Returns:
        操作结果字符串
    """
    try:
        prs = _safe_open_ppt(file_path)
        slide_layout = prs.slide_layouts[0]  # 标题布局
        slide = prs.slides.add_slide(slide_layout)

        # 设置背景
        if bg_color:
            _set_slide_bg_color(slide, bg_color)

        # 设置标题
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = title
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    if title_color:
                        parsed = _parse_color(title_color)
                        if parsed:
                            run.font.color.rgb = parsed
                    run.font.size = Pt(44)
                    run.font.bold = True

        # 设置副标题
        if subtitle:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:
                    shape.text = subtitle
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if subtitle_color:
                                parsed = _parse_color(subtitle_color)
                                if parsed:
                                    run.font.color.rgb = parsed
                            run.font.size = Pt(24)
                    break

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            total = len(prs.slides)
            return f"✓ 标题幻灯片已添加，共 {total} 张幻灯片"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 添加标题幻灯片失败: {str(e)}"


def add_bullet_slide(file_path: str, title: str = "内容",
                     bullets: list = None,
                     font_name: str = None,
                     font_size: int = 18,
                     color: str = None,
                     bg_color: str = None) -> str:
    """
    添加一张带项目符号的内容幻灯片。

    Args:
        file_path: PPT 文件路径
        title: 幻灯片标题
        bullets: 项目符号列表，如 ['要点1', '要点2', '要点3']
        font_name: 字体名称
        font_size: 字号（磅）
        color: 文字颜色
        bg_color: 背景颜色

    Returns:
        操作结果字符串
    """
    if not bullets:
        bullets = []

    try:
        prs = _safe_open_ppt(file_path)
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)

        # 背景
        if bg_color:
            _set_slide_bg_color(slide, bg_color)

        # 标题
        if slide.shapes.title:
            slide.shapes.title.text = title

        # 项目符号内容
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                tf = shape.text_frame
                tf.clear()

                for i, bullet_text in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = str(bullet_text)
                    p.level = 0

                    for run in p.runs:
                        if font_name:
                            run.font.name = font_name
                        if font_size:
                            run.font.size = Pt(font_size)
                        if color:
                            parsed = _parse_color(color)
                            if parsed:
                                run.font.color.rgb = parsed

                break

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            total = len(prs.slides)
            return f"✓ 项目符号幻灯片已添加，共 {len(bullets)} 条要点，共 {total} 张幻灯片"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 添加内容幻灯片失败: {str(e)}"


def read_ppt_info(file_path: str) -> str:
    """
    读取 PPT 文件的基本信息。

    Args:
        file_path: PPT 文件路径

    Returns:
        包含幻灯片数量、尺寸等信息的字符串
    """
    if not os.path.exists(file_path):
        return f"[ERROR] 文件 '{file_path}' 不存在。"

    try:
        prs = Presentation(file_path)
        info_lines = [
            f"文件: {file_path}",
            f"幻灯片数量: {len(prs.slides)}",
            f"幻灯片宽度: {prs.slide_width.inches:.2f} 英寸 ({prs.slide_width.emu} EMU)",
            f"幻灯片高度: {prs.slide_height.inches:.2f} 英寸 ({prs.slide_height.emu} EMU)",
            "",
        ]

        for i, slide in enumerate(prs.slides, 1):
            shape_count = len(slide.shapes)
            layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
            info_lines.append(f"--- 幻灯片 {i} (布局: {layout_name}, {shape_count} 个形状) ---")

            for shape in slide.shapes:
                shape_info = f"  - {shape.shape_type}: "
                if hasattr(shape, 'text') and shape.text:
                    text_preview = shape.text[:50].replace("\n", " ")
                    shape_info += f"'{text_preview}'"
                else:
                    shape_info += f"位置=({shape.left}, {shape.top}), 尺寸=({shape.width}, {shape.height})"
                info_lines.append(shape_info)

        return "\n".join(info_lines)
    except Exception as e:
        return f"[ERROR] 读取 PPT 信息失败: {str(e)}"


def read_ppt_text(file_path: str) -> str:
    """
    提取 PPT 文件中所有幻灯片的文字内容。

    Args:
        file_path: PPT 文件路径

    Returns:
        包含所有文字内容的字符串
    """
    if not os.path.exists(file_path):
        return f"[ERROR] 文件 '{file_path}' 不存在。"

    try:
        prs = Presentation(file_path)
        output = [f"文件 '{file_path}' 的文字内容:"]

        for i, slide in enumerate(prs.slides, 1):
            output.append(f"\n--- 幻灯片 {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    output.append(f"  {shape.text.strip()}")

        return "\n".join(output)
    except Exception as e:
        return f"[ERROR] 读取 PPT 文字失败: {str(e)}"


def duplicate_slide(file_path: str, slide_index: int = 1) -> str:
    """
    复制指定幻灯片并追加到末尾。

    Args:
        file_path: PPT 文件路径
        slide_index: 要复制的幻灯片索引（从 1 开始）

    Returns:
        操作结果字符串
    """
    if not os.path.exists(file_path):
        return f"[ERROR] 文件 '{file_path}' 不存在。"

    try:
        prs = Presentation(file_path)
        slides = prs.slides
        if slide_index < 1 or slide_index > len(slides):
            return f"[ERROR] 幻灯片索引 {slide_index} 超出范围 (1-{len(slides)})。"

        # 获取源幻灯片
        source_slide = slides[slide_index - 1]

        # 复制幻灯片（通过底层 XML）
        slide_layout = source_slide.slide_layout
        new_slide = prs.slides.add_slide(slide_layout)

        # 复制形状
        for shape in source_slide.shapes:
            el = shape._element
            new_el = el.__deepcopy__(None)

            # 添加形状到新幻灯片
            spTree = new_slide.shapes._spTree
            spTree.append(new_el)

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            total = len(prs.slides)
            return f"✓ 已复制第 {slide_index} 张幻灯片，共 {total} 张幻灯片"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 复制幻灯片失败: {str(e)}"


def delete_slide(file_path: str, slide_index: int = 1) -> str:
    """
    删除指定幻灯片。

    Args:
        file_path: PPT 文件路径
        slide_index: 要删除的幻灯片索引（从 1 开始）

    Returns:
        操作结果字符串
    """
    if not os.path.exists(file_path):
        return f"[ERROR] 文件 '{file_path}' 不存在。"

    try:
        prs = Presentation(file_path)
        slides = prs.slides
        if slide_index < 1 or slide_index > len(slides):
            return f"[ERROR] 幻灯片索引 {slide_index} 超出范围 (1-{len(slides)})。"

        # 通过底层 XML 删除幻灯片
        sldIdLst = prs.presentation.sldIdLst
        slide_id = slides[slide_index - 1].slide_id
        for sldId in sldIdLst:
            if sldId.get(qn('r:id')) == slide_id or int(sldId.get('id')) == slide_id:
                sldIdLst.remove(sldId)
                break

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            total = len(prs.slides) if hasattr(prs, 'slides') else 0
            return f"✓ 已删除第 {slide_index} 张幻灯片，剩余 {total} 张"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 删除幻灯片失败: {str(e)}"


def reorder_slides(file_path: str, new_order: list = None) -> str:
    """
    重新排列幻灯片顺序。

    Args:
        file_path: PPT 文件路径
        new_order: 新的顺序列表，如 [3, 1, 2] 表示将第3张移至第1位

    Returns:
        操作结果字符串
    """
    if not os.path.exists(file_path):
        return f"[ERROR] 文件 '{file_path}' 不存在。"

    try:
        prs = Presentation(file_path)
        slides = list(prs.slides)
        n = len(slides)

        if not new_order:
            return "[ERROR] 请提供新的幻灯片顺序。"
        if len(new_order) != n:
            return f"[ERROR] 新顺序长度 ({len(new_order)}) 与幻灯片数量 ({n}) 不匹配。"
        if sorted(new_order) != list(range(1, n + 1)):
            return "[ERROR] 新顺序必须包含 1 到 N 的所有整数。"

        # 通过操作底层 XML 重排
        sldIdLst = prs.presentation.sldIdLst
        sldId_elements = list(sldIdLst)

        # 按新顺序重新排列
        reordered = [sldId_elements[i - 1] for i in new_order]

        # 清空并重新添加
        for element in sldId_elements:
            sldIdLst.remove(element)
        for element in reordered:
            sldIdLst.append(element)

        saved = _safe_save_ppt(prs, file_path)
        if saved:
            return f"✓ 幻灯片顺序已调整为: {new_order}"
        else:
            return f"[ERROR] 保存 PPT 文件失败"
    except Exception as e:
        return f"[ERROR] 重排幻灯片失败: {str(e)}"
